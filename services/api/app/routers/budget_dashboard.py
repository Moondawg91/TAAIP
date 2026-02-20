from fastapi import APIRouter, Request, HTTPException
from fastapi import Depends
from typing import Dict, Any, List
from ..db import connect, row_to_dict
from fastapi.responses import Response
from .rbac import get_current_user, require_any_role

router = APIRouter(prefix="/budget", tags=["budget"])


def _safe_sum(cur, sql, params=()):
    try:
        cur.execute(sql, params)
        r = cur.fetchone()
        if not r:
            return 0.0
        # r may be dict-like
        val = list(r.values())[0] if hasattr(r, 'values') else (r[0] if isinstance(r, (list, tuple)) else r)
        return float(val) if val is not None else 0.0
    except Exception:
        return 0.0


@router.get('/dashboard')
def budget_dashboard(request: Request, fy: int = None, qtr: int = None, org_unit_id: int = None, station_id: str = None, funding_line: str = None, funding_source: str = None, eor_code: str = None) -> Dict[str, Any]:
    conn = connect()
    try:
        cur = conn.cursor()

        filters = {}
        where_parts = []
        params: List[Any] = []
        # For budget_line_item we join to fy_budget to filter by fy/org_unit
        fy_clause = ''
        if fy is not None:
            filters['fy'] = fy
            fy_clause = ' AND fb.fy=?'
            params.append(fy)
        if qtr is not None:
            filters['qtr'] = qtr
        if org_unit_id is not None:
            filters['org_unit_id'] = org_unit_id
            if 'fy' not in filters:
                # still need to join fy_budget if filtering by org_unit
                pass
        if station_id is not None:
            filters['station_id'] = station_id
        if funding_line is not None:
            filters['funding_line'] = funding_line
        if funding_source is not None:
            filters['funding_source'] = funding_source
        # permission gate optional national advertising funding
        if funding_source == 'ADVERTISING_FUNDS_NATIONAL':
            # Require RBAC role/permission rather than a custom header.
            # We call the RBAC helpers directly for a best-effort check in-process.
            try:
                user = get_current_user(request)
                dep = require_any_role('advertising_access')
                # call dependency function directly with resolved user
                dep(user)
            except HTTPException:
                raise
            except Exception:
                # any unexpected error -> forbid
                raise HTTPException(status_code=403, detail='forbidden_funding_source')
        if eor_code is not None:
            filters['eor_code'] = eor_code

        # allocated: sum budget_line_item.amount joined to fy_budget for filter by fy/org
        # allocated: prefer amounts tied to a fy_budget (fb) which has fy and org_unit
        alloc_params = []
        sql_alloc = 'SELECT SUM(COALESCE(bli.amount,0)) as s FROM budget_line_item bli LEFT JOIN fy_budget fb ON fb.id = bli.fy_budget_id WHERE 1=1'
        if fy is not None:
            sql_alloc += ' AND fb.fy=?'; alloc_params.append(fy)
        if org_unit_id is not None:
            sql_alloc += ' AND fb.org_unit_id=?'; alloc_params.append(org_unit_id)
        if funding_source is not None:
            sql_alloc += ' AND bli.funding_source=?'; alloc_params.append(funding_source)
        if eor_code is not None:
            sql_alloc += ' AND bli.eor_code=?'; alloc_params.append(eor_code)
        allocated = _safe_sum(cur, sql_alloc, tuple(alloc_params))

        # fallback: include any budget_line_item rows that are not linked to a fy_budget
        # (some imports may write standalone line items). Those rows don't have fy/org
        # metadata in this schema, so we sum amounts where fy_budget_id IS NULL.
        try:
            fsql = 'SELECT SUM(COALESCE(amount,0)) as s FROM budget_line_item WHERE fy_budget_id IS NULL'
            unlinked = _safe_sum(cur, fsql)
            allocated = allocated + unlinked
        except Exception:
            pass

        # planned: sum of projects.planned_cost + event.planned_cost
        try:
            sql_proj = 'SELECT SUM(COALESCE(planned_cost,0)) as s FROM projects WHERE 1=1'
            pparams = []
            if fy is not None:
                sql_proj += ' AND fy=?'; pparams.append(fy)
            if org_unit_id is not None:
                sql_proj += ' AND org_unit_id=?'; pparams.append(org_unit_id)
            if funding_line is not None:
                sql_proj += ' AND funding_line=?'; pparams.append(funding_line)
            if funding_source is not None:
                sql_proj += ' AND funding_source=?'; pparams.append(funding_source)
            if eor_code is not None:
                sql_proj += ' AND eor_code=?'; pparams.append(eor_code)
            proj_planned = _safe_sum(cur, sql_proj, tuple(pparams))
        except Exception:
            proj_planned = 0.0
        try:
            sql_evt = 'SELECT SUM(COALESCE(planned_cost,0)) as s FROM event WHERE 1=1'
            eparams = []
            if fy is not None:
                sql_evt += ' AND fy=?'; eparams.append(fy)
            if org_unit_id is not None:
                sql_evt += ' AND org_unit_id=?'; eparams.append(org_unit_id)
            if funding_line is not None:
                sql_evt += ' AND funding_line=?'; eparams.append(funding_line)
            if funding_source is not None:
                sql_evt += ' AND funding_source=?'; eparams.append(funding_source)
            if eor_code is not None:
                sql_evt += ' AND eor_code=?'; eparams.append(eor_code)
            evt_planned = _safe_sum(cur, sql_evt, tuple(eparams))
        except Exception:
            evt_planned = 0.0
        planned = proj_planned + evt_planned

        # actual: sum of expenses.amount
        try:
            sql_act = 'SELECT SUM(COALESCE(amount,0)) as s FROM expenses WHERE 1=1'
            aparams = []
            if fy is not None:
                sql_act += ' AND fy=?'; aparams.append(fy)
            if org_unit_id is not None:
                sql_act += ' AND org_unit_id=?'; aparams.append(org_unit_id)
            if funding_line is not None:
                sql_act += ' AND funding_line=?'; aparams.append(funding_line)
            if funding_source is not None:
                sql_act += ' AND funding_source=?'; aparams.append(funding_source)
            if eor_code is not None:
                sql_act += ' AND eor_code=?'; aparams.append(eor_code)
            actual = _safe_sum(cur, sql_act, tuple(aparams))
        except Exception:
            actual = 0.0

        # Compute obligated: budget line items with committed/obligated statuses
        try:
            sql_ob = 'SELECT SUM(COALESCE(bli.amount,0)) as s FROM budget_line_item bli LEFT JOIN fy_budget fb ON fb.id = bli.fy_budget_id WHERE (bli.status IN ("committed","obligated","committed" ) OR bli.status IS NOT NULL)'
            ob_params = []
            if fy is not None:
                sql_ob += ' AND fb.fy=?'; ob_params.append(fy)
            if org_unit_id is not None:
                sql_ob += ' AND fb.org_unit_id=?'; ob_params.append(org_unit_id)
            if funding_source is not None:
                sql_ob += ' AND bli.funding_source=?'; ob_params.append(funding_source)
            if eor_code is not None:
                sql_ob += ' AND bli.eor_code=?'; ob_params.append(eor_code)
            obligated = _safe_sum(cur, sql_ob, tuple(ob_params))
        except Exception:
            obligated = 0.0

        executed = actual
        remaining = allocated - obligated - executed

        # by_category: union budget_line_item.category, projects.category, event.category, expenses.category
        by_category = []
        try:
            sql_cat = '''
            SELECT category as category, SUM(allocated) as allocated, SUM(planned) as planned, SUM(actual) as actual
            FROM (
                SELECT bli.category as category, COALESCE(bli.amount,0) as allocated, 0 as planned, 0 as actual FROM budget_line_item bli LEFT JOIN fy_budget fb ON fb.id = bli.fy_budget_id
                UNION ALL
                SELECT category, 0 as allocated, COALESCE(planned_cost,0) as planned, 0 as actual FROM projects
                UNION ALL
                SELECT category, 0, COALESCE(planned_cost,0), 0 FROM event
                UNION ALL
                SELECT category, 0, 0, COALESCE(amount,0) FROM expenses
            ) GROUP BY category
            '''
            cur.execute(sql_cat)
            rows = cur.fetchall()
            for r in rows:
                cat = r['category'] if r and 'category' in r and r['category'] is not None else 'uncategorized'
                a = float(r['allocated'] or 0)
                p = float(r['planned'] or 0)
                ac = float(r['actual'] or 0)
                by_category.append({'category': cat, 'allocated': a, 'planned': p, 'actual': ac, 'remaining': a - p - ac})
        except Exception:
            by_category = []

        # by_funding_line
        by_funding_line = []
        try:
            sql_fl = '''
            SELECT funding_line as funding_line, SUM(allocated) as allocated, SUM(planned) as planned, SUM(actual) as actual
            FROM (
                SELECT fb.funding_line as funding_line, COALESCE(bli.amount,0) as allocated, 0 as planned, 0 as actual FROM budget_line_item bli LEFT JOIN fy_budget fb ON fb.id = bli.fy_budget_id
                UNION ALL
                SELECT funding_line, 0, COALESCE(planned_cost,0), 0 FROM projects
                UNION ALL
                SELECT funding_line, 0, COALESCE(planned_cost,0), 0 FROM event
                UNION ALL
                SELECT funding_line, 0, 0, COALESCE(amount,0) FROM expenses
            ) GROUP BY funding_line
            '''
            cur.execute(sql_fl)
            for r in cur.fetchall():
                fl = r['funding_line'] if r and 'funding_line' in r and r['funding_line'] is not None else 'unassigned'
                a = float(r['allocated'] or 0); p = float(r['planned'] or 0); ac = float(r['actual'] or 0)
                by_funding_line.append({'funding_line': fl, 'allocated': a, 'obligated': 0.0, 'executed': ac, 'remaining': a - ac})
        except Exception:
            by_funding_line = []

        # by_project: list projects with planned and actual from expenses
        by_project = []
        try:
            cur.execute('SELECT project_id as project_id, title as name, COALESCE(planned_cost,0) as planned FROM projects')
            prows = cur.fetchall()
            for pr in prows:
                pid = pr['project_id'] if 'project_id' in pr else pr.get('id')
                name = pr.get('name') or pr.get('title') or ''
                planned_p = float(pr.get('planned') or 0)
                cur.execute('SELECT SUM(COALESCE(amount,0)) as s FROM expenses WHERE project_id=?', (pid,))
                rr = cur.fetchone(); actual_p = float(rr['s']) if rr and rr['s'] is not None else 0.0
            by_project.append({'project_id': pid, 'name': name, 'allocated': 0.0, 'obligated': 0.0, 'executed': actual_p, 'remaining': max(planned_p - actual_p, 0.0)})
        except Exception:
            by_project = []

        # by_event
        by_event = []
        try:
            cur.execute('SELECT id as event_id, name, COALESCE(planned_cost,0) as planned, project_id FROM event')
            erows = cur.fetchall()
            for er in erows:
                eid = er['event_id']
                name = er.get('name')
                planned_e = float(er.get('planned') or 0)
                proj_id = er.get('project_id')
                cur.execute('SELECT SUM(COALESCE(amount,0)) as s FROM expenses WHERE event_id=?', (eid,))
                rr = cur.fetchone(); actual_e = float(rr['s']) if rr and rr['s'] is not None else 0.0
                by_event.append({'event_id': eid, 'name': name, 'project_id': proj_id, 'planned': planned_e, 'actual': actual_e, 'variance': planned_e - actual_e})
        except Exception:
            by_event = []

        # missing_data hints
        missing = []
        try:
            cur.execute('SELECT COUNT(1) as c FROM budget_line_item')
            r = cur.fetchone(); if_budget = int(r['c']) if r and r['c'] is not None else 0
            if if_budget == 0:
                missing.append('No budget lines imported')
        except Exception:
            missing.append('No budget lines table')
        try:
            cur.execute('SELECT COUNT(1) as c FROM projects')
            r = cur.fetchone(); if_proj = int(r['c']) if r and r['c'] is not None else 0
            if if_proj == 0:
                missing.append('No projects imported')
        except Exception:
            missing.append('No projects table')
        try:
            cur.execute('SELECT COUNT(1) as c FROM expenses')
            r = cur.fetchone(); if_exp = int(r['c']) if r and r['c'] is not None else 0
            if if_exp == 0:
                missing.append('No expenses imported')
        except Exception:
            missing.append('No expenses table')

        # Return Phase-10 shaped payload
        totals = {
            'allocated': round(allocated, 2),
            'obligated': round(obligated, 2),
            'executed': round(executed, 2),
            'remaining': round(remaining, 2)
        }

        # breakdown by funding_source
        by_funding_source = []
        try:
            # sum allocated/obligated/executed grouped by funding_source
            sql_bfs = "SELECT COALESCE(bli.funding_source, 'unassigned') as funding_source, SUM(COALESCE(bli.amount,0)) as allocated FROM budget_line_item bli LEFT JOIN fy_budget fb ON fb.id=bli.fy_budget_id WHERE 1=1"
            bfs_params = []
            if fy is not None:
                sql_bfs += ' AND fb.fy=?'; bfs_params.append(fy)
            if org_unit_id is not None:
                sql_bfs += ' AND fb.org_unit_id=?'; bfs_params.append(org_unit_id)
            if funding_source is not None:
                sql_bfs += ' AND bli.funding_source=?'; bfs_params.append(funding_source)
            if eor_code is not None:
                sql_bfs += ' AND bli.eor_code=?'; bfs_params.append(eor_code)
            sql_bfs += ' GROUP BY bli.funding_source'
            cur.execute(sql_bfs, tuple(bfs_params))
            for r in cur.fetchall():
                fs = r.get('funding_source') or 'unassigned'
                a = float(r.get('allocated') or 0)
                # executed for this funding_source from expenses
                cur.execute('SELECT SUM(COALESCE(amount,0)) as s FROM expenses WHERE funding_source=?', (fs,))
                rr = cur.fetchone(); exec_fs = float(rr['s']) if rr and rr['s'] is not None else 0.0
                by_funding_source.append({'funding_source': fs, 'allocated': round(a,2), 'obligated': 0.0, 'executed': round(exec_fs,2), 'remaining': round(a - exec_fs,2)})
        except Exception:
            by_funding_source = []

        # breakdown by EOR
        by_eor = []
        try:
            sql_be = "SELECT COALESCE(bli.eor_code, 'unassigned') as eor_code, SUM(COALESCE(bli.amount,0)) as allocated FROM budget_line_item bli LEFT JOIN fy_budget fb ON fb.id=bli.fy_budget_id WHERE 1=1"
            be_params = []
            if fy is not None:
                sql_be += ' AND fb.fy=?'; be_params.append(fy)
            if org_unit_id is not None:
                sql_be += ' AND fb.org_unit_id=?'; be_params.append(org_unit_id)
            if funding_source is not None:
                sql_be += ' AND bli.funding_source=?'; be_params.append(funding_source)
            if eor_code is not None:
                sql_be += ' AND bli.eor_code=?'; be_params.append(eor_code)
            sql_be += ' GROUP BY bli.eor_code'
            cur.execute(sql_be, tuple(be_params))
            for r in cur.fetchall():
                e = r.get('eor_code') or 'unassigned'
                a = float(r.get('allocated') or 0)
                cur.execute('SELECT SUM(COALESCE(amount,0)) as s FROM expenses WHERE eor_code=?', (e,))
                rr = cur.fetchone(); exec_e = float(rr['s']) if rr and rr['s'] is not None else 0.0
                by_eor.append({'eor_code': e, 'allocated': round(a,2), 'obligated': 0.0, 'executed': round(exec_e,2), 'remaining': round(a - exec_e,2)})
        except Exception:
            by_eor = []

        # Ensure deterministic numeric zeros
        def zero_safe_list(lst, key_fields):
            out = []
            for item in lst:
                o = {}
                for k in key_fields:
                    o[k] = item.get(k) if item.get(k) is not None else 0
                out.append(item)
            return out

        payload = {
            'filters': filters,
            'totals': totals,
            'by_funding_source': by_funding_source,
            'by_eor': by_eor,
            'by_project': by_project,
            'by_event': by_event,
            'by_category': by_category,
            'missing_data': missing
        }

        # Backwards-compatible top-level KPI keys expected by older callers/tests
        payload['total_planned'] = round(planned, 2)
        payload['total_spent'] = round(executed, 2)
        payload['total_pending'] = round(obligated, 2)
        payload['total_remaining'] = round(remaining, 2)

        # New Phase-10 shaped `kpis` object for dashboard callers
        payload['kpis'] = {
            'allocated': round(allocated, 2),
            'planned': round(planned, 2),
            'actual': round(executed, 2),
            'remaining': round((allocated - planned - executed), 2)
        }

        # Backwards-compatible breakdown keys
        payload['breakdown_by_category'] = by_category
        payload['breakdown_by_funding_line'] = by_funding_line

        return payload
    finally:
        try:
            conn.close()
        except Exception:
            pass


@router.get('/dashboard/export.csv')
def budget_dashboard_export(request: Request, fy: int = None, qtr: int = None, org_unit_id: int = None, station_id: str = None, funding_line: str = None, funding_source: str = None, eor_code: str = None):
    # produce a simple CSV with KPI and breakdown rows
    import csv, io
    data = budget_dashboard(request, fy=fy, qtr=qtr, org_unit_id=org_unit_id, station_id=station_id, funding_line=funding_line, funding_source=funding_source, eor_code=eor_code)
    # normalize KPI shape
    kpis = {
        'total_planned': data.get('total_planned'),
        'total_spent': data.get('total_spent'),
        'total_pending': data.get('total_pending'),
        'total_remaining': data.get('total_remaining')
    }
    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow(['metric','value'])
    for k, v in kpis.items():
        writer.writerow([k, v])
    writer.writerow([])
    writer.writerow(['category','allocated','planned','actual','remaining'])
    for c in (data.get('breakdown_by_category') or []):
        writer.writerow([c.get('category'), c.get('allocated'), c.get('planned'), c.get('actual'), c.get('remaining')])
    writer.writerow([])
    writer.writerow(['funding_line','allocated','planned','actual','remaining'])
    for f in (data.get('breakdown_by_funding_line') or []):
        writer.writerow([f.get('funding_line'), f.get('allocated'), f.get('planned'), f.get('actual'), f.get('remaining')])
    csv_text = buf.getvalue()
    return Response(content=csv_text, media_type='text/csv', headers={
        'Content-Disposition': 'attachment; filename="budget_dashboard.csv"'
    })


@router.get('/dashboard/export.json')
def budget_dashboard_export_json(fy: int = None, qtr: int = None, org_unit_id: int = None, station_id: str = None, funding_line: str = None, funding_source: str = None, eor_code: str = None):
    data = budget_dashboard(fy=fy, qtr=qtr, org_unit_id=org_unit_id, station_id=station_id, funding_line=funding_line, funding_source=funding_source, eor_code=eor_code)
    # return JSON
    from fastapi.responses import JSONResponse
    return JSONResponse(content=data)


@router.get('/dashboard/export.xlsx')
def budget_dashboard_export_xlsx(fy: int = None, qtr: int = None, org_unit_id: int = None, station_id: str = None, funding_line: str = None, funding_source: str = None, eor_code: str = None):
    data = budget_dashboard(fy=fy, qtr=qtr, org_unit_id=org_unit_id, station_id=station_id, funding_line=funding_line, funding_source=funding_source, eor_code=eor_code)
    try:
        import openpyxl
        from openpyxl.workbook import Workbook
        from openpyxl.utils import get_column_letter
        wb = Workbook()
        ws = wb.active
        ws.title = 'KPIs'
        kpis = [
            ('total_planned', data.get('total_planned')),
            ('total_spent', data.get('total_spent')),
            ('total_pending', data.get('total_pending')),
            ('total_remaining', data.get('total_remaining')),
        ]
        ws.append(['metric','value'])
        for row in kpis:
            ws.append(row)
        # add category sheet
        ws2 = wb.create_sheet('By Category')
        ws2.append(['category','allocated','planned','actual','remaining'])
        for c in (data.get('breakdown_by_category') or []):
            ws2.append([c.get('category'), c.get('allocated'), c.get('planned'), c.get('actual'), c.get('remaining')])
        # write to bytes
        from io import BytesIO
        bio = BytesIO()
        wb.save(bio)
        bio.seek(0)
        return Response(content=bio.read(), media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', headers={'Content-Disposition': 'attachment; filename="budget_dashboard.xlsx"'})
    except Exception:
        # fallback to CSV
        return budget_dashboard_export(fy=fy, qtr=qtr, org_unit_id=org_unit_id, station_id=station_id, funding_line=funding_line)


@router.get('/dashboard/export.pdf')
def budget_dashboard_export_pdf(fy: int = None, qtr: int = None, org_unit_id: int = None, station_id: str = None, funding_line: str = None, funding_source: str = None, eor_code: str = None):
    data = budget_dashboard(fy=fy, qtr=qtr, org_unit_id=org_unit_id, station_id=station_id, funding_line=funding_line, funding_source=funding_source, eor_code=eor_code)
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        from io import BytesIO
        bio = BytesIO()
        c = canvas.Canvas(bio, pagesize=letter)
        y = 750
        c.setFont('Helvetica-Bold', 12)
        c.drawString(50, y, 'Budget Dashboard KPIs')
        y -= 20
        c.setFont('Helvetica', 10)
        for k in ['total_planned','total_spent','total_pending','total_remaining']:
            c.drawString(50, y, f"{k}: {data.get(k)}")
            y -= 16
        c.showPage()
        c.save()
        bio.seek(0)
        return Response(content=bio.read(), media_type='application/pdf', headers={'Content-Disposition': 'attachment; filename="budget_dashboard.pdf"'})
    except Exception:
        # fallback to CSV
        return budget_dashboard_export(fy=fy, qtr=qtr, org_unit_id=org_unit_id, station_id=station_id, funding_line=funding_line)


@router.get('/dashboard/export.pptx')
def budget_dashboard_export_pptx(fy: int = None, qtr: int = None, org_unit_id: int = None, station_id: str = None, funding_line: str = None, funding_source: str = None, eor_code: str = None):
    data = budget_dashboard(fy=fy, qtr=qtr, org_unit_id=org_unit_id, station_id=station_id, funding_line=funding_line, funding_source=funding_source, eor_code=eor_code)
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from io import BytesIO
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = 'Budget Dashboard KPIs'
        left = Inches(1)
        top = Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, Inches(8), Inches(4))
        tf = txBox.text_frame
        for k in ['total_planned','total_spent','total_pending','total_remaining']:
            p = tf.add_paragraph()
            p.text = f"{k}: {data.get(k)}"
        bio = BytesIO()
        prs.save(bio)
        bio.seek(0)
        return Response(content=bio.read(), media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', headers={'Content-Disposition': 'attachment; filename="budget_dashboard.pptx"'})
    except Exception:
        # fallback to CSV
        return budget_dashboard_export(fy=fy, qtr=qtr, org_unit_id=org_unit_id, station_id=station_id, funding_line=funding_line)
